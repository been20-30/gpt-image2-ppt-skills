from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[2] / "outputs/arabic-courseware-8/demo-arabic.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BG = RGBColor(20, 19, 18); WHITE = RGBColor(247, 246, 242); BLUE = RGBColor(55, 78, 232); LIME = RGBColor(190, 255, 48); MUTED = RGBColor(185, 184, 180)

def box(slide, x,y,w,h, text="", size=22, color=WHITE, bold=False, align=PP_ALIGN.RIGHT, font="Noto Kufi Arabic"):
    sh=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=sh.text_frame; tf.clear(); tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=align; p.space_after=Pt(0); r=p.add_run(); r.text=text; r.font.name=font; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return sh

def base(title, kicker=None):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg=s.background.fill; bg.solid(); bg.fore_color.rgb=BG
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.8), Inches(.35), Inches(.12), Inches(6.7)); bar.fill.solid(); bar.fill.fore_color.rgb=BLUE; bar.line.fill.background()
    if kicker: box(s, 9.0,.45,3.2,.35,kicker,11,LIME,True)
    box(s, 5.5,.65,6.5,.7,title,26,WHITE,True)
    return s

# 1 hero
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
box(s, 5.0,1.55,7.4,1.0,"التعلم الذكي يبدأ بسؤال جيد",34,WHITE,True)
box(s, 5.2,2.75,7.2,.8,"كيف نستخدم الذكاء الاصطناعي لتوسيع التفكير، لا لاستبداله؟",18,WHITE)
box(s, .9,5.8,5.3,.45,"دليل عملي للطلاب والمهنيين",14,LIME,True)
# 2 editorial
s=base("المشكلة ليست نقص المعلومات","01 · CONTEXT")
box(s, 1.25,2.0,5.7,1.5,"وفرة المعرفة لا تضمن جودة القرار",28,WHITE,True)
box(s, 1.3,3.8,5.3,1.0,"التحدي الحقيقي هو اختيار السؤال، وتحديد معيار الإجابة، وتحويل المعرفة إلى قرار قابل للتنفيذ.",17,MUTED)
# 3 framework
s=base("إطار التعلم الذكي","02 · FRAMEWORK")
for i,(t,c) in enumerate([("اسأل",BLUE),("اختبر",LIME),("طبّق",BLUE),("راجع",LIME)]):
    x=1.3+i*2.75; sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.6), Inches(2.15), Inches(1.2)); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background(); box(s,x+.1,2.85,1.95,.45,t,20,BG,True)
box(s,1.3,4.5,9.8,.65,"كل دورة قصيرة تزيد جودة الفهم وتحوّل المعرفة إلى ممارسة.",17,MUTED)
# 4 timeline
s=base("من السؤال إلى التجربة","03 · PROCESS")
for i,(n,t) in enumerate([("01","سؤال محدد"),("02","مثال مضاد"),("03","تجربة صغيرة"),("04","مراجعة النتيجة")]):
    y=2.0+i*1.05; box(s,1.25,y,1.0,.45,n,14,LIME,True,PP_ALIGN.LEFT,"IBM Plex Sans"); box(s,2.3,y,3.8,.45,t,19,WHITE,True); line=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.4), Inches(y+.2), Inches(4.8), Inches(.02)); line.fill.solid(); line.fill.fore_color.rgb=BLUE; line.line.fill.background()
# 5 comparison
s=base("إجابة سريعة أم فهم عميق؟","04 · COMPARISON")
box(s,1.25,2.1,4.6,.6,"إجابة سريعة",22,BLUE,True); box(s,7.0,2.1,4.6,.6,"فهم عميق",22,LIME,True)
box(s,1.25,3.0,4.5,1.6,"تمنحك بداية\nلكنها قد تخفي الافتراضات",19,WHITE,True)
box(s,7.0,3.0,4.5,1.6,"يطلب دليلًا\nويوضح الحدود ويصنع ممارسة",19,WHITE,True)
# 6 big number
s=base("ثلاث عادات تصنع فرقًا","05 · TAKEAWAY")
box(s,1.3,2.0,2.4,2.0,"3",72,LIME,True,PP_ALIGN.LEFT,"IBM Plex Sans")
box(s,4.0,2.15,6.5,.6,"عادات عملية",28,WHITE,True)
box(s,4.0,3.1,6.8,1.3,"اكتب ما تعرفه قبل أن تسأل\nاطلب تفسيرًا لا نتيجة فقط\nاحتفظ بسجل قصير لما تعلّمته",18,MUTED)
# 7 quote
s=base("لا تقِس التعلم بعدد الإجابات","06 · INSIGHT")
box(s,1.4,2.0,9.8,1.7,"«التعلم الحقيقي هو أن يصبح سؤالك التالي أدق من سؤالك السابق.»",27,WHITE,True)
box(s,1.4,4.5,8.0,.55,"الفكرة ليست جمع مخرجات أكثر، بل بناء حكم أفضل.",18,LIME,True)
# 8 closing
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
box(s,5.0,1.7,7.2,.8,"ابدأ بتجربة واحدة اليوم",32,WHITE,True)
box(s,5.2,3.0,6.8,1.0,"اختر سؤالًا حقيقيًا، وطبّق دورة التعلم الذكي مرة واحدة.",20,MUTED)
box(s,1.0,5.9,6.4,.45,"الخطوة التالية: حوّل الفضول إلى ممارسة.",16,LIME,True)
prs.save(OUT); print(OUT)
