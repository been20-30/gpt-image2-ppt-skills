"""Render an EditableScene as native PowerPoint objects."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .scene import EditableScene, SceneElement


ALIGNMENTS = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


def _hex_color(value: str) -> RGBColor:
    normalized = value.strip().lstrip("#")
    if len(normalized) != 6:
        raise ValueError(f"无效颜色: {value}")
    return RGBColor.from_string(normalized.upper())


def _slide_bbox(prs: Presentation, scene: EditableScene, element: SceneElement) -> tuple[int, int, int, int]:
    x, y, width, height = element.bbox_px
    return (
        int(x / scene.canvas_width * prs.slide_width),
        int(y / scene.canvas_height * prs.slide_height),
        int(width / scene.canvas_width * prs.slide_width),
        int(height / scene.canvas_height * prs.slide_height),
    )


def _add_text(prs: Presentation, slide, scene: EditableScene, element: SceneElement) -> None:
    left, top, width, height = _slide_bbox(prs, scene, element)
    shape = slide.shapes.add_textbox(left, top, width, height)
    shape.name = element.id
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    paragraph = frame.paragraphs[0]
    paragraph.alignment = ALIGNMENTS.get(str(element.style.get("align", "left")), PP_ALIGN.LEFT)
    run = paragraph.add_run()
    run.text = element.content
    font = run.font
    font.name = str(element.style.get("font_face", "Noto Sans CJK SC"))
    font.size = Pt(float(element.style.get("font_size_pt", 24)))
    font.bold = int(element.style.get("font_weight", 400)) >= 600
    font.color.rgb = _hex_color(str(element.style.get("color", "#000000")))


def render_editable_pptx(scene: EditableScene, output_path: Path | str) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    clean_plate = slide.shapes.add_picture(
        str(scene.clean_plate),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    clean_plate.name = "clean_plate"

    for element in sorted(scene.elements, key=lambda item: item.z_index):
        if element.type == "image_layer":
            left, top, width, height = _slide_bbox(prs, scene, element)
            picture = slide.shapes.add_picture(str(element.asset), left, top, width=width, height=height)
            picture.name = element.id
        elif element.type == "native_text":
            _add_text(prs, slide, scene, element)

    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    return path
